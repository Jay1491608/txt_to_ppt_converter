"""
Theme-adaptive TXT-to-PPTX converter.

Converts structured text into a PowerPoint deck that inherits styling from a
template file — no hardcoded colors, fonts, or positions.
"""

from __future__ import annotations

import argparse
import logging
import re
import uuid
from pathlib import Path
from typing import Any

from lxml.etree import Element
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)

_P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
_SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"


def _p14(tag: str) -> str:
    return f"{{{_P14_NS}}}{tag}"

# Layout name fragments to search for (order matters for fallback preference)
TOPIC_DIVIDER_LAYOUT_NAMES = ("title slide",)
SECTION_LAYOUT_NAMES = ("section header", "section title", "divider")
CONTENT_LAYOUT_NAMES = ("title and content", "content with caption")
TWO_COLUMN_LAYOUT_NAMES = ("two content", "comparison")

SLIDE_HEADING_RE = re.compile(r"^### Slide \d+ — (.+)$")
SLIDE_M89_RE = re.compile(r"^## SLIDE\s*(\d+)\s*—\s*(.+)$", re.IGNORECASE)
PLAINTEXT_SLIDE_RE = re.compile(r"^SLIDE\s*(\d+)\s*\|\s*Topic\s*(\d+):\s*(.+)$", re.IGNORECASE)
PLAINTEXT_SECTION_RE = re.compile(r"^SECTION SLIDE\s*—\s*TOPIC\s*(\d+):\s*(.+)$", re.IGNORECASE)
PLAINTEXT_MODULE_RE = re.compile(r"^MODULE\s*(\d+)\s*—\s*(.+)$", re.IGNORECASE)
PLAINTEXT_TITLE_RE = re.compile(r"^Title:\s*(.+)$", re.IGNORECASE)
PLAINTEXT_BULLET_RE = re.compile(r"^(\s*)[•\-\*]\s*(.+)$")
PLAINTEXT_QA_Q_RE = re.compile(r"^Q(\d+):\s*(.+)$", re.IGNORECASE)
PLAINTEXT_QA_A_RE = re.compile(r"^A(\d+):\s*(.+)$", re.IGNORECASE)
PLAINTEXT_SEPARATOR_RE = re.compile(r"^([=\-#])\1+$")
DOC_TITLE_RE = re.compile(r"^# (.+)$")
TOPIC_CODE_RE = re.compile(r"^## TOPIC (\d+)\.(\d+)\b")
TOPIC_INLINE_RE = re.compile(r"^TOPIC (\d+):\s*(.+)$", re.IGNORECASE)
TOPIC_HEADING_RE = re.compile(r"^## TOPIC\b")
MODULE_FILE_RE = re.compile(r"^M(\d+)$", re.IGNORECASE)
MODULE_LINE_RE = re.compile(r"^\*{0,2}Module\s*(\d+)\s*[:—\-]\s*(.+?)\*{0,2}$", re.IGNORECASE)
BULLET_RE = re.compile(r"^(\s*)- (.+)$")
QA_ITEM_RE = re.compile(r"^\d+\.\s+\*\*Q:\*\*")
QA_SECTION_RE = re.compile(r"^Q\s*&\s*A(\s+Section)?:?\s*$", re.IGNORECASE)
QA_NUM_RE = re.compile(r"^Q(\d+):\s*(.+)$", re.IGNORECASE)
QA_ANSWER_RE = re.compile(r"^A:\s*(.+)$", re.IGNORECASE)
SEPARATOR_LINE_RE = re.compile(r"^=+$")
FOOTER_LINE_RE = re.compile(
    r"^This completes all \d+ slides for Module \d+",
    re.IGNORECASE,
)
SUMMARY_LINE_RE = re.compile(
    r"^(# END OF|PRESENTATION SUMMARY|Total:\s*\d+\s*Slides|\d+\.\s*Topic \d+:)",
    re.IGNORECASE,
)


def _strip_inline_markdown(text: str) -> str:
    """Remove simple **bold** markers while keeping the label text."""
    return re.sub(r"\*\*(.+?)\*\*", r"\1", text).strip()


def _detect_format(path: Path) -> str:
    with path.open(encoding="utf-8") as handle:
        for raw_line in handle:
            stripped = raw_line.strip()
            if SLIDE_HEADING_RE.match(stripped):
                return "m1"
            if SLIDE_M89_RE.match(stripped):
                return "m89"
            if PLAINTEXT_SLIDE_RE.match(stripped) or PLAINTEXT_SECTION_RE.match(stripped):
                return "plaintext"
    return "simple"


M1_FIELD_PREFIXES = (
    "**Title:**",
    "**Bullets:**",
    "**Speaker Notes:**",
    "**Q&A:**",
)


def parse_m1_txt_to_structure(filepath: str | Path) -> dict[str, Any]:
    """
    Parse Module-style structured text (e.g. M1.txt) into slides.

    Convention:
      # Document Title                    -> presentation metadata (not a slide)
      ## TOPIC X.X — Topic Name           -> topic header (not a slide)
      ### Slide N — Section Slide         -> section divider slide
      **Title:** Section name             -> title for the section slide above
      ### Slide N — Content Title         -> content slide
      **Bullets:**                        -> starts bullet list for current slide
      - bullet text                       -> slide bullet (supports **bold** labels)
      **Speaker Notes:** text             -> speaker notes for current slide
      **Q&A:**                            -> Q&A block appended to speaker notes
      1. **Q:** ... **A:** ...            -> individual Q&A items
    """
    structure: dict[str, Any] = {"title": "", "module": None, "module_code": "", "sections": []}
    current: dict[str, Any] | None = None
    current_topic_code: str | None = None
    mode: str | None = None
    path = Path(filepath)

    with path.open(encoding="utf-8") as handle:
        for raw_line in handle:
            line = raw_line.rstrip("\n\r")
            stripped = line.strip()

            if not stripped or stripped == "---" or SEPARATOR_LINE_RE.match(stripped):
                continue

            if FOOTER_LINE_RE.match(stripped) or SUMMARY_LINE_RE.match(stripped):
                continue

            if stripped.startswith("- ") and SUMMARY_LINE_RE.search(
                structure.get("title", "")
            ):
                # Skip summary bullet lists appended at end of some module files.
                continue

            doc_title_match = DOC_TITLE_RE.match(stripped)
            if doc_title_match:
                title_text = doc_title_match.group(1).strip()
                if not structure["title"]:
                    structure["title"] = title_text
                mod_match = re.search(r"\bModule\s*(\d+)\b", title_text, re.IGNORECASE)
                mod_title_match = re.search(r"\bModule\s*\d+\s*[:—\-]\s*(.+)", title_text, re.IGNORECASE)
                if mod_match:
                    m_num = int(mod_match.group(1))
                    m_code = f"M{m_num}"
                    structure["module"] = m_num
                    structure["module_code"] = m_code
                else:
                    file_match = MODULE_FILE_RE.match(path.stem)
                    m_code = f"M{file_match.group(1)}" if file_match else (structure.get("module_code") or "")

                slide_title = mod_title_match.group(1).strip() if mod_title_match else title_text

                current = {
                    "type": "module_slide",
                    "title": slide_title,
                    "bullets": [],
                    "notes": "",
                    "topic_code": m_code,
                    "module_code": m_code,
                }
                structure["sections"].append(current)
                continue

            module_line_match = MODULE_LINE_RE.match(stripped)
            if module_line_match:
                m_num = int(module_line_match.group(1))
                m_title = module_line_match.group(2).strip()
                m_code = f"M{m_num}"
                structure["module"] = m_num
                structure["module_code"] = m_code

                if current and current.get("type") == "module_slide":
                    current["title"] = m_title
                    current["topic_code"] = m_code
                    current["module_code"] = m_code
                else:
                    current = {
                        "type": "module_slide",
                        "title": m_title,
                        "bullets": [],
                        "notes": "",
                        "topic_code": m_code,
                        "module_code": m_code,
                    }
                    structure["sections"].append(current)
                continue

            topic_match = TOPIC_CODE_RE.match(stripped)
            if topic_match:
                module_num, topic_num = topic_match.groups()
                structure["module"] = int(module_num)
                structure["module_code"] = f"M{module_num}"
                current_topic_code = f"M{module_num}T{topic_num}"
                continue

            inline_topic_match = TOPIC_INLINE_RE.match(stripped)
            if inline_topic_match:
                topic_num = inline_topic_match.group(1)
                module_num = structure.get("module")
                if module_num is None:
                    file_match = MODULE_FILE_RE.match(path.stem)
                    if file_match:
                        module_num = int(file_match.group(1))
                        structure["module"] = module_num
                        structure["module_code"] = f"M{module_num}"
                if module_num is not None:
                    current_topic_code = f"M{module_num}T{topic_num}"
                continue

            if TOPIC_HEADING_RE.match(stripped) or (
                stripped.startswith("**")
                and stripped.endswith("**")
                and not any(stripped.startswith(prefix) for prefix in M1_FIELD_PREFIXES)
            ):
                continue

            slide_match = SLIDE_HEADING_RE.match(stripped)
            if slide_match:
                slide_label = slide_match.group(1).strip()
                if slide_label.lower() == "section slide":
                    current = {
                        "type": "section_divider",
                        "title": "",
                        "bullets": [],
                        "notes": "",
                        "topic_code": current_topic_code or "",
                        "module_code": structure.get("module_code") or "",
                    }
                else:
                    current = {
                        "type": "content",
                        "title": slide_label,
                        "bullets": [],
                        "notes": "",
                        "topic_code": current_topic_code or "",
                        "module_code": structure.get("module_code") or "",
                    }
                structure["sections"].append(current)
                mode = None
                continue

            if current is None:
                logger.warning("Unrecognized line (no active slide): %r", stripped)
                continue

            if stripped.startswith("**Title:**"):
                current["title"] = _strip_inline_markdown(stripped[len("**Title:**") :])
                continue

            if stripped == "**Bullets:**":
                mode = "bullets"
                continue

            if stripped.startswith("**Speaker Notes:**"):
                note_text = _strip_inline_markdown(stripped[len("**Speaker Notes:**") :])
                current["notes"] = (
                    f"{current['notes']}\n\n{note_text}".strip()
                    if current["notes"]
                    else note_text
                )
                mode = None
                continue

            if stripped == "**Q&A:**" or QA_SECTION_RE.match(stripped):
                mode = "qa"
                if current["notes"]:
                    current["notes"] = f"{current['notes']}\n\nQ&A:"
                else:
                    current["notes"] = "Q&A:"
                continue

            if mode == "qa" and QA_ITEM_RE.match(stripped):
                qa_line = _strip_inline_markdown(stripped)
                current["notes"] = f"{current['notes']}\n{qa_line}"
                continue

            if mode == "qa":
                q_match = QA_NUM_RE.match(stripped)
                if q_match:
                    qa_line = f"Q{q_match.group(1)}: {q_match.group(2).strip()}"
                    current["notes"] = f"{current['notes']}\n{qa_line}"
                    continue
                a_match = QA_ANSWER_RE.match(stripped)
                if a_match:
                    current["notes"] = f"{current['notes']}\nA: {a_match.group(1).strip()}"
                    continue

            bullet_match = BULLET_RE.match(line)
            if bullet_match and mode in ("bullets", None):
                indent = len(bullet_match.group(1))
                level = 1 if indent >= 2 else 0
                current["bullets"].append(
                    {
                        "text": _strip_inline_markdown(bullet_match.group(2)),
                        "level": level,
                    }
                )
                mode = "bullets"
                continue

            logger.warning("Unrecognized line in slide %r: %r", current.get("title"), stripped)

    if not structure["sections"]:
        raise ValueError(f"No slides found in {path}. Expected '### Slide N — ...' headings.")

    for item in structure["sections"]:
        if item["type"] == "section_divider" and not item["title"]:
            raise ValueError(
                "Section slide is missing **Title:** immediately after '### Slide N — Section Slide'."
            )

    _infer_module_from_path(path, structure)
    _assign_section_names(structure)
    return structure


def _parse_plaintext_to_structure(filepath: Path) -> dict[str, Any]:
    """
    Parse plaintext structured text (e.g. M2.txt, M3.txt) into slides.
    """
    structure: dict[str, Any] = {"title": "", "module": None, "module_code": "", "sections": []}
    current: dict[str, Any] | None = None
    current_topic_code: str | None = None
    mode: str | None = None
    path = Path(filepath)

    _infer_module_from_path(path, structure)
    module_code = structure.get("module_code") or ""

    with path.open(encoding="utf-8") as handle:
        for raw_line in handle:
            line = raw_line.rstrip("\n\r")
            stripped = line.strip()

            if not stripped or PLAINTEXT_SEPARATOR_RE.match(stripped):
                continue

            if stripped.startswith("(") and stripped.endswith(")"):
                continue

            if (stripped.startswith("PowerPoint Content Pack") or 
                "Content Pack" in stripped or
                "slides (" in stripped or
                stripped.startswith("Note:") or
                "AI for Government Professionals" in stripped):
                continue

            module_match = PLAINTEXT_MODULE_RE.match(stripped)
            if module_match:
                module_num = int(module_match.group(1))
                structure["module"] = module_num
                structure["module_code"] = f"M{module_num}"
                module_code = f"M{module_num}"
                title_text = module_match.group(2).strip()
                if not structure["title"]:
                    structure["title"] = title_text
                current = {
                    "type": "module_slide",
                    "title": f"MODULE {module_num} — {title_text}",
                    "bullets": [],
                    "notes": "",
                    "topic_code": module_code,
                    "module_code": module_code,
                }
                structure["sections"].append(current)
                continue

            section_match = PLAINTEXT_SECTION_RE.match(stripped)
            if section_match:
                topic_num = section_match.group(1)
                topic_title = section_match.group(2).strip()
                current_topic_code = f"{module_code}T{topic_num}" if module_code else f"T{topic_num}"
                current = {
                    "type": "section_divider",
                    "title": topic_title,
                    "bullets": [],
                    "notes": "",
                    "topic_code": current_topic_code,
                    "module_code": module_code,
                }
                structure["sections"].append(current)
                mode = None
                continue

            slide_match = PLAINTEXT_SLIDE_RE.match(stripped)
            if slide_match:
                slide_num = slide_match.group(1)
                topic_num = slide_match.group(2)
                current_topic_code = f"{module_code}T{topic_num}" if module_code else f"T{topic_num}"
                current = {
                    "type": "content",
                    "title": "",
                    "bullets": [],
                    "notes": "",
                    "topic_code": current_topic_code,
                    "module_code": module_code,
                }
                structure["sections"].append(current)
                mode = None
                continue

            if current is None:
                continue

            title_match = PLAINTEXT_TITLE_RE.match(stripped)
            if title_match:
                current["title"] = title_match.group(1).strip()
                continue

            if stripped.upper() == "BULLET POINTS:":
                mode = "bullets"
                continue

            if stripped.upper() == "SPEAKER NOTES:":
                mode = "notes"
                continue

            if stripped.upper() == "Q&A:":
                mode = "qa"
                if current["notes"]:
                    current["notes"] = f"{current['notes']}\n\nQ&A:"
                else:
                    current["notes"] = "Q&A:"
                continue

            bullet_match = PLAINTEXT_BULLET_RE.match(line)
            if bullet_match and mode == "bullets":
                indent = len(bullet_match.group(1))
                level = 1 if indent >= 2 else 0
                current["bullets"].append({
                    "text": bullet_match.group(2).strip(),
                    "level": level
                })
                continue

            if mode == "notes":
                if current["notes"]:
                    current["notes"] = f"{current['notes']}\n{stripped}"
                else:
                    current["notes"] = stripped
                continue

            if mode == "qa":
                q_match = PLAINTEXT_QA_Q_RE.match(stripped)
                if q_match:
                    qa_line = f"Q{q_match.group(1)}: {q_match.group(2).strip()}"
                    current["notes"] = f"{current['notes']}\n{qa_line}"
                    continue
                a_match = PLAINTEXT_QA_A_RE.match(stripped)
                if a_match:
                    qa_line = f"A{a_match.group(1)}: {a_match.group(2).strip()}"
                    current["notes"] = f"{current['notes']}\n{qa_line}"
                    continue

    _assign_section_names(structure)
    return structure


def _parse_m89_to_structure(filepath: Path) -> dict[str, Any]:
    """
    Parse Module 8/9 markdown format (e.g. M8.txt, M9.txt) into slides.
    """
    structure: dict[str, Any] = {"title": "", "module": None, "module_code": "", "sections": []}
    current: dict[str, Any] | None = None
    current_topic_code: str | None = None
    mode: str | None = None
    path = Path(filepath)

    _infer_module_from_path(path, structure)
    module_code = structure.get("module_code") or ""

    with path.open(encoding="utf-8") as handle:
        for raw_line in handle:
            line = raw_line.rstrip("\n\r")
            stripped = line.strip()

            if not stripped or stripped == "---" or SEPARATOR_LINE_RE.match(stripped):
                continue

            doc_title_match = DOC_TITLE_RE.match(stripped)
            if doc_title_match:
                title_text = doc_title_match.group(1).strip()
                if not structure["title"]:
                    structure["title"] = title_text
                mod_match = re.search(r"\bModule\s*(\d+)\b", title_text, re.IGNORECASE)
                if mod_match:
                    m_num = int(mod_match.group(1))
                    m_code = f"M{m_num}"
                    structure["module"] = m_num
                    structure["module_code"] = m_code
                    module_code = m_code
                current = {
                    "type": "module_slide",
                    "title": title_text,
                    "bullets": [],
                    "notes": "",
                    "topic_code": module_code,
                    "module_code": module_code,
                }
                structure["sections"].append(current)
                continue

            slide_match = SLIDE_M89_RE.match(stripped)
            if slide_match:
                slide_num = slide_match.group(1)
                slide_label = slide_match.group(2).strip()

                if slide_label.lower() == "section title":
                    current = {
                        "type": "section_divider",
                        "title": "",
                        "bullets": [],
                        "notes": "",
                        "topic_code": current_topic_code or "",
                        "module_code": module_code,
                    }
                else:
                    current = {
                        "type": "content",
                        "title": slide_label,
                        "bullets": [],
                        "notes": "",
                        "topic_code": current_topic_code or "",
                        "module_code": module_code,
                    }
                structure["sections"].append(current)
                mode = None
                continue

            if current is None:
                continue

            if current["type"] == "section_divider" and not current["title"] and stripped.startswith("**") and stripped.endswith("**"):
                bold_text = _strip_inline_markdown(stripped)
                current["title"] = bold_text

                topic_match = re.match(r"^(\d+)\.(\d+)\b", bold_text)
                if topic_match:
                    m_num, t_num = topic_match.groups()
                    structure["module"] = int(m_num)
                    structure["module_code"] = f"M{m_num}"
                    module_code = f"M{m_num}"
                    current_topic_code = f"M{m_num}T{t_num}"
                    current["topic_code"] = current_topic_code
                    current["module_code"] = module_code
                continue

            if stripped.startswith("**Speaker Notes:**") or stripped.startswith("Speaker Notes:"):
                mode = "notes"
                prefix = "**Speaker Notes:**" if stripped.startswith("**Speaker Notes:**") else "Speaker Notes:"
                note_text = _strip_inline_markdown(stripped[len(prefix) :])
                current["notes"] = note_text
                continue

            if stripped.startswith("**Q&A:**") or stripped.startswith("Q&A:"):
                mode = "qa"
                if current["notes"]:
                    current["notes"] = f"{current['notes']}\n\nQ&A:"
                else:
                    current["notes"] = "Q&A:"
                continue

            bullet_match = BULLET_RE.match(line)
            if bullet_match and mode in ("bullets", None):
                indent = len(bullet_match.group(1))
                level = 1 if indent >= 2 else 0
                current["bullets"].append({
                    "text": _strip_inline_markdown(bullet_match.group(2)),
                    "level": level
                })
                mode = "bullets"
                continue

            if mode == "notes":
                current["notes"] = f"{current['notes']}\n{stripped}"
                continue

            if mode == "qa":
                qa_match = re.match(r"^(\d+)\.\s*Q:\s*(.+?)\s+A:\s*(.+)$", stripped, re.IGNORECASE)
                if qa_match:
                    q_num = qa_match.group(1)
                    q_text = qa_match.group(2).strip()
                    a_text = qa_match.group(3).strip()
                    current["notes"] = f"{current['notes']}\nQ{q_num}: {q_text}\nA{q_num}: {a_text}"
                else:
                    current["notes"] = f"{current['notes']}\n{stripped}"
                continue

    for item in structure["sections"]:
        if item["type"] == "section_divider" and not item["title"]:
            item["title"] = "Section Divider"

    _assign_section_names(structure)
    return structure


def _assign_section_names(structure: dict[str, Any]) -> None:
    """
    Set section_name on each slide based on per-module topic layout.

    - Modules with multiple topics (e.g. M1: M1T1–M1T10) → topic sections
    - Modules with a single topic (e.g. M2–M9) → module section (M2, M3, …)
    """
    slides = structure.get("sections", [])
    topics_by_module: dict[str, set[str]] = {}

    for item in slides:
        module_code = item.get("module_code", "")
        topic_code = item.get("topic_code", "")
        if module_code and topic_code and topic_code != module_code:
            topics_by_module.setdefault(module_code, set()).add(topic_code)

    for item in slides:
        module_code = item.get("module_code", "")
        topic_code = item.get("topic_code", "")
        topic_count = len(topics_by_module.get(module_code, set()))

        if topic_count > 1 and topic_code:
            item["section_name"] = topic_code
        elif module_code:
            item["section_name"] = module_code
        else:
            item["section_name"] = topic_code or "Default Section"


def _infer_module_from_path(path: Path, structure: dict[str, Any]) -> None:
    """Set module/module_code from filename (e.g. M1.txt) when not parsed from TOPIC headings."""
    file_match = MODULE_FILE_RE.match(path.stem)
    if file_match and structure.get("module") is None:
        module_num = int(file_match.group(1))
        structure["module"] = module_num
        structure["module_code"] = f"M{module_num}"

    module_code = structure.get("module_code") or ""
    for item in structure["sections"]:
        if not item.get("module_code"):
            item["module_code"] = module_code
        if not item.get("topic_code") and module_code:
            item["topic_code"] = module_code


def merge_structures(structures: list[dict[str, Any]]) -> dict[str, Any]:
    """Combine multiple parsed modules into one deck structure."""
    if not structures:
        raise ValueError("No structures to merge.")

    merged: dict[str, Any] = {
        "title": structures[0].get("title", ""),
        "module": None,
        "module_code": "",
        "sections": [],
    }
    for structure in structures:
        merged["sections"].extend(structure["sections"])
    _assign_section_names(merged)
    return merged


def parse_txt_to_structure(filepath: str | Path) -> dict[str, Any]:
    """Parse input text, auto-detecting simple or Module (M1) format."""
    path = Path(filepath)
    fmt = _detect_format(path)
    if fmt == "m1":
        logger.info("Detected Module-style input format (### Slide headings)")
        return parse_m1_txt_to_structure(path)
    elif fmt == "m89":
        logger.info("Detected Module 8/9 markdown format (## SLIDE N headings)")
        return _parse_m89_to_structure(path)
    elif fmt == "plaintext":
        logger.info("Detected Plaintext structured format (SLIDE N | Topic headings)")
        return _parse_plaintext_to_structure(path)

    return _parse_simple_txt_to_structure(path)


def _parse_simple_txt_to_structure(filepath: Path) -> dict[str, Any]:
    """
    Parse a structured text file into an intermediate slide model.

    Convention:
      # Section Title       -> new section divider slide
      ## Slide Title        -> new content slide
      - bullet text         -> top-level bullet
        - sub bullet text   -> nested bullet (2-space indent before '-')
      Notes: ...            -> speaker notes for the current slide

    Blank lines are ignored. Lines before the first heading are ignored.
    """
    structure: dict[str, Any] = {"title": "", "sections": []}
    current: dict[str, Any] | None = None
    path = Path(filepath)

    with path.open(encoding="utf-8") as handle:
        for raw_line in handle:
            line = raw_line.rstrip("\n\r")
            if not line.strip():
                continue

            if line.startswith("# "):
                current = {
                    "type": "section_divider",
                    "title": line[2:].strip(),
                    "bullets": [],
                    "notes": "",
                }
                structure["sections"].append(current)
                if not structure["title"]:
                    structure["title"] = current["title"]
                continue

            if line.startswith("## "):
                current = {
                    "type": "content",
                    "title": line[3:].strip(),
                    "bullets": [],
                    "notes": "",
                }
                structure["sections"].append(current)
                continue

            if line.startswith("Notes:") and current is not None:
                note_text = line[len("Notes:") :].strip()
                if current["notes"]:
                    current["notes"] = f"{current['notes']}\n{note_text}"
                else:
                    current["notes"] = note_text
                continue

            bullet_match = re.match(r"^(\s*)- (.+)$", line)
            if bullet_match and current is not None:
                indent = len(bullet_match.group(1))
                level = 1 if indent >= 2 else 0
                current["bullets"].append(
                    {"text": bullet_match.group(2).strip(), "level": level}
                )
                continue

            logger.warning("Unrecognized line (no active slide): %r", line)

    if not structure["sections"]:
        raise ValueError(
            f"No slides found in {path}. Use '# Section' or '## Slide Title' headings."
        )

    return structure


def _placeholder_label(placeholder: Any) -> str:
    try:
        ph_type = placeholder.placeholder_format.type
        return PP_PLACEHOLDER(ph_type).name
    except (ValueError, AttributeError):
        return "UNKNOWN"


def validate_template(prs: Presentation) -> None:
    """Log layout names and placeholder indices so template mismatches fail fast."""
    logger.info("Template validation — %d slide layout(s):", len(prs.slide_layouts))
    for index, layout in enumerate(prs.slide_layouts):
        placeholders = list(layout.placeholders)
        if placeholders:
            details = ", ".join(
                f"idx={ph.placeholder_format.idx} ({_placeholder_label(ph)})"
                for ph in placeholders
            )
        else:
            details = "no placeholders"
        logger.info("  [%d] %r — %s", index, layout.name, details)


def find_layout(
    prs: Presentation,
    name_fragments: tuple[str, ...],
    fallback_index: int = 1,
) -> Any:
    """Return the first slide layout whose name contains any of the fragments."""
    lowered = tuple(fragment.lower() for fragment in name_fragments)
    for layout in prs.slide_layouts:
        layout_name = layout.name.lower()
        if any(fragment in layout_name for fragment in lowered):
            logger.debug("Matched layout %r for fragments %s", layout.name, name_fragments)
            return layout

    if fallback_index < len(prs.slide_layouts):
        fallback = prs.slide_layouts[fallback_index]
        logger.warning(
            "No layout matched %s; falling back to [%d] %r",
            name_fragments,
            fallback_index,
            fallback.name,
        )
        return fallback

    raise ValueError(
        f"Could not find a layout matching {name_fragments} "
        f"and fallback index {fallback_index} is out of range."
    )


def _find_title_placeholder(slide: Any) -> Any | None:
    if slide.shapes.title is not None:
        return slide.shapes.title
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            return placeholder
    return None


def _find_body_placeholder(slide: Any) -> Any | None:
    preferred_types = (
        PP_PLACEHOLDER.BODY,
        PP_PLACEHOLDER.OBJECT,
    )
    for ph_type in preferred_types:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == ph_type:
                return placeholder

    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx != 0:
            return placeholder
    return None


def _fill_bullets(text_frame: Any, bullets: list[dict[str, Any]]) -> None:
    text_frame.clear()
    first = True
    for bullet in bullets:
        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        paragraph.text = bullet["text"]
        paragraph.level = bullet["level"]
        first = False


def _clear_existing_slides(prs: Presentation) -> int:
    """Remove all slides from a presentation while keeping masters, layouts, and theme."""
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(slide_id)
    return len(slide_ids)


def _remove_existing_powerpoint_sections(prs: Presentation) -> None:
    ext_lists = prs.element.findall(qn("p:extLst"))
    if not ext_lists:
        return
    ext_list = ext_lists[0]
    for ext in list(ext_list.findall(qn("p:ext"))):
        if ext.get("uri") == _SECTION_EXT_URI:
            ext_list.remove(ext)


def apply_powerpoint_sections(
    prs: Presentation,
    section_groups: list[tuple[str, list[Any]]],
) -> None:
    """Assign PowerPoint section names (e.g. M1T1) to groups of slides."""
    if not section_groups:
        return

    _remove_existing_powerpoint_sections(prs)

    ext_lists = prs.element.findall(qn("p:extLst"))
    if ext_lists:
        ext_list = ext_lists[0]
    else:
        ext_list = Element(qn("p:extLst"))
        prs.element.append(ext_list)

    ext = Element(qn("p:ext"))
    ext.set("uri", _SECTION_EXT_URI)
    ext_list.append(ext)

    section_list = Element(_p14("sectionLst"))
    ext.append(section_list)

    for section_name, slides in section_groups:
        if not slides:
            continue

        section = Element(_p14("section"))
        section.set("name", section_name)
        section.set("id", f"{{{uuid.uuid4()}}}")
        section_list.append(section)

        slide_id_list = Element(_p14("sldIdLst"))
        section.append(slide_id_list)

        for slide in slides:
            slide_id = Element(_p14("sldId"))
            slide_id.set("id", str(slide.slide_id))
            slide_id_list.append(slide_id)

        logger.info("PowerPoint section %r: %d slide(s)", section_name, len(slides))


def _section_name_for_slide(item: dict[str, Any]) -> str:
    return (
        item.get("section_name")
        or item.get("topic_code")
        or item.get("module_code")
        or "Default Section"
    )


def _section_sort_key(name: str) -> tuple[int, str]:
    match = re.search(r"M(\d+)", name)
    return (int(match.group(1)), name) if match else (999, name)


def _log_section_summary(structure: dict[str, Any]) -> None:
    """Log slide counts per PowerPoint section and module."""
    section_counts: dict[str, int] = {}
    module_counts: dict[str, int] = {}
    for item in structure["sections"]:
        section_name = _section_name_for_slide(item)
        section_counts[section_name] = section_counts.get(section_name, 0) + 1
        module_code = item.get("module_code", "")
        if module_code:
            module_counts[module_code] = module_counts.get(module_code, 0) + 1

    logger.info("Slide plan — %d total slide(s) from content:", len(structure["sections"]))
    for module_code in sorted(module_counts, key=lambda name: int(name[1:])):
        logger.info("  %s: %d slide(s)", module_code, module_counts[module_code])
    logger.info("PowerPoint sections:")
    for section_name in sorted(section_counts, key=_section_sort_key):
        logger.info("  %s: %d slide(s)", section_name, section_counts[section_name])


def _divider_layout_for_item(prs: Presentation, item: dict[str, Any]) -> Any:
    if item.get("topic_code"):
        return find_layout(prs, TOPIC_DIVIDER_LAYOUT_NAMES, fallback_index=0)
    return find_layout(prs, SECTION_LAYOUT_NAMES, fallback_index=2)


def build_pptx_from_structure(
    structure: dict[str, Any],
    template_path: str | Path,
    output_path: str | Path,
    *,
    clear_template_slides: bool = True,
    apply_sections: bool = True,
) -> None:
    """Build a PPTX deck from parsed structure using template layouts and placeholders."""
    template_path = Path(template_path)
    output_path = Path(output_path)

    prs = Presentation(str(template_path))
    validate_template(prs)
    _log_section_summary(structure)

    existing_slides = len(prs.slides)
    if existing_slides and clear_template_slides:
        removed = _clear_existing_slides(prs)
        logger.info(
            "Template used for format only — cleared %d reference slide(s), kept theme/layouts",
            removed,
        )
    elif existing_slides:
        logger.warning(
            "Template contains %d existing slide(s); new slides will be appended",
            existing_slides,
        )

    content_layout = find_layout(prs, CONTENT_LAYOUT_NAMES, fallback_index=1)

    section_groups: list[tuple[str, list[Any]]] = []
    current_section_name: str | None = None
    current_section_slides: list[Any] = []

    title_layout = find_layout(prs, TOPIC_DIVIDER_LAYOUT_NAMES, fallback_index=0)

    for slide_index, item in enumerate(structure["sections"], start=1):
        if item["type"] in ("module_slide", "section_divider"):
            layout = title_layout
        else:
            layout = content_layout
        slide = prs.slides.add_slide(layout)

        section_name = _section_name_for_slide(item)
        if section_name != current_section_name:
            if current_section_slides:
                section_groups.append((current_section_name or "Default Section", current_section_slides))
            current_section_name = section_name
            current_section_slides = []
        current_section_slides.append(slide)

        title_placeholder = _find_title_placeholder(slide)
        if title_placeholder is None:
            logger.error(
                "Slide %d (%r): no title placeholder on layout %r",
                slide_index,
                item["title"],
                layout.name,
            )
        else:
            title_placeholder.text = item["title"]
            logger.info(
                "Slide %d: type=%s, layout=%r, title placeholder idx=%s",
                slide_index,
                item["type"],
                layout.name,
                title_placeholder.placeholder_format.idx,
            )

        if item["bullets"]:
            body_placeholder = _find_body_placeholder(slide)
            if body_placeholder is None:
                logger.error(
                    "Slide %d (%r): no body placeholder for %d bullet(s) on layout %r",
                    slide_index,
                    item["title"],
                    len(item["bullets"]),
                    layout.name,
                )
            else:
                _fill_bullets(body_placeholder.text_frame, item["bullets"])
                logger.info(
                    "Slide %d: body placeholder idx=%s (%s), %d bullet(s)",
                    slide_index,
                    body_placeholder.placeholder_format.idx,
                    _placeholder_label(body_placeholder),
                    len(item["bullets"]),
                )

        if item["notes"]:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = item["notes"]
            logger.info("Slide %d: speaker notes set (%d chars)", slide_index, len(item["notes"]))

    if current_section_slides:
        section_groups.append((current_section_name or "Default Section", current_section_slides))

    if apply_sections and section_groups:
        apply_powerpoint_sections(prs, section_groups)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info("Saved %d slide(s) to %s", len(structure["sections"]), output_path)


def convert_txt_to_pptx(
    input_txt: str | Path,
    template_pptx: str | Path,
    output_pptx: str | Path,
    *,
    clear_template_slides: bool = True,
    apply_sections: bool = True,
) -> dict[str, Any]:
    structure = parse_txt_to_structure(input_txt)
    logger.info(
        "Parsed %d slide(s) from %s (presentation title: %r, module: %s)",
        len(structure["sections"]),
        input_txt,
        structure.get("title") or "(none)",
        structure.get("module_code") or "(none)",
    )
    build_pptx_from_structure(
        structure,
        template_pptx,
        output_pptx,
        clear_template_slides=clear_template_slides,
        apply_sections=apply_sections,
    )
    return structure


def convert_modules_batch(
    input_files: list[str | Path],
    template_pptx: str | Path,
    output_dir: str | Path,
    *,
    combine: bool = False,
    combined_output: str | Path | None = None,
    clear_template_slides: bool = True,
) -> list[Path]:
    """Convert M1.txt … M9.txt into separate PPTX files or one combined deck."""
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    sorted_files = sorted(Path(f) for f in input_files)
    structures: list[dict[str, Any]] = []
    outputs: list[Path] = []

    for input_path in sorted_files:
        structure = parse_txt_to_structure(input_path)
        structures.append(structure)

        if combine:
            continue

        module_code = structure.get("module_code") or input_path.stem
        output_path = output_dir / f"{module_code}.pptx"
        build_pptx_from_structure(
            structure,
            template_pptx,
            output_path,
            clear_template_slides=clear_template_slides,
        )
        logger.info("Wrote %s (%d slides)", output_path, len(structure["sections"]))
        outputs.append(output_path)

    if combine:
        if not combined_output:
            combined_output = output_dir / "full_course.pptx"
        merged = merge_structures(structures)
        build_pptx_from_structure(
            merged,
            template_pptx,
            combined_output,
            clear_template_slides=clear_template_slides,
        )
        logger.info(
            "Wrote combined deck %s (%d slides from %d module(s))",
            combined_output,
            len(merged["sections"]),
            len(structures),
        )
        outputs.append(Path(combined_output))

    return outputs


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Convert structured TXT/MD content into a theme-adaptive PPTX deck."
    )
    parser.add_argument(
        "input_txt",
        nargs="?",
        help="Path to a single structured .txt or .md input file",
    )
    parser.add_argument(
        "-t",
        "--template",
        dest="template_flag",
        help="Format-reference template (-t/--template or positional); slides are not kept",
    )
    parser.add_argument(
        "template_pptx",
        nargs="?",
        help="Format-reference .pptx/.potx (theme and layouts only; slides are replaced)",
    )
    parser.add_argument(
        "output_pptx",
        nargs="?",
        help="Path for the generated .pptx output (single-file mode)",
    )
    parser.add_argument(
        "--batch",
        nargs="+",
        metavar="TXT",
        help="Convert multiple module files (e.g. examples/M1.txt examples/M2.txt)",
    )
    parser.add_argument(
        "--combine",
        action="store_true",
        help="With --batch, merge all modules into one PPTX instead of separate files",
    )
    parser.add_argument(
        "--output-dir",
        default="output",
        help="Output folder for --batch mode (default: output)",
    )
    parser.add_argument(
        "-o",
        "--combined-output",
        help="Output path for --batch --combine (default: output/full_course.pptx)",
    )
    parser.add_argument(
        "-v",
        "--verbose",
        action="store_true",
        help="Enable debug logging (layout matching, placeholder details)",
    )
    parser.add_argument(
        "--keep-template-slides",
        action="store_true",
        help="Keep existing slides in the template file (default: clear them first)",
    )
    parser.add_argument(
        "--no-sections",
        action="store_true",
        help="Do not create PowerPoint section groups (M1T1, M1T2, ...)",
    )
    args = parser.parse_args(argv)

    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(levelname)s: %(message)s",
    )

    clear_template_slides = not args.keep_template_slides
    apply_sections = not args.no_sections
    template_path = args.template_flag or args.template_pptx

    if args.batch:
        if not template_path:
            parser.error("--batch requires a template (-t/--template or positional template_pptx)")
        convert_modules_batch(
            args.batch,
            template_path,
            args.output_dir,
            combine=args.combine,
            combined_output=args.combined_output,
            clear_template_slides=clear_template_slides,
        )
        return 0

    if not args.input_txt or not template_path or not args.output_pptx:
        parser.error("single-file mode requires: input_txt template_pptx output_pptx")

    convert_txt_to_pptx(
        args.input_txt,
        template_path,
        args.output_pptx,
        clear_template_slides=clear_template_slides,
        apply_sections=apply_sections,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
