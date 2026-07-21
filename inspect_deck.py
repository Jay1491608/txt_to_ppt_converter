"""Inspect slide titles/sections in the template deck."""
from pptx import Presentation

prs = Presentation(r"output/AI for Energy Proffessional.pptx")
print(f"Total slides: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides, 1):
    title = ""
    if slide.shapes.title and slide.shapes.title.text:
        title = slide.shapes.title.text.strip()
    else:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0 and hasattr(shape, "text"):
                title = shape.text.strip()
                break
    layout = slide.slide_layout.name
    # show first 50 and any with M*T* pattern
    if i <= 60 or "M" in title[:4]:
        print(f"{i:3d} | {layout:20s} | {title[:80]}")
