"""Inspect template layouts."""
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

prs = Presentation(r"output/AI for Energy Proffessional.pptx")
print("Slides in template:", len(prs.slides))
print("Layouts:")
for i, layout in enumerate(prs.slide_layouts):
    phs = []
    for ph in layout.placeholders:
        try:
            t = PP_PLACEHOLDER(ph.placeholder_format.type).name
        except ValueError:
            t = str(ph.placeholder_format.type)
        phs.append(f"idx={ph.placeholder_format.idx}({t})")
    joined = ", ".join(phs) if phs else "no placeholders"
    print(f"  [{i}] {layout.name!r} -> {joined}")
