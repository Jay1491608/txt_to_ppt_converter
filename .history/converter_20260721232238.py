"""
Theme-adaptive TXT-to-PPTX converter.

Converts structured text into a PowerPoint deck that inherits styling from a
template file — no hardcoded colors, fonts, or positions.
"""

from __future__ import annotations

import argparse
import logging
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

logger = logging.getLogger(__name__)

# Layout name fragments to search for (order matters for fallback preference)
SECTION_LAYOUT_NAMES = ("section header", "section title", "divider")
CONTENT_LAYOUT_NAMES = ("title and content", "content with caption")
TWO_COLUMN_LAYOUT_NAMES = ("two content", "comparison")

SLIDE_HEADING_RE = re.compile(r"^### Slide \d+ — (.+)$")
DOC_TITLE_RE = re.compile(r"^# (.+)$")
TOPIC_HEADING_RE = re.compile(r"^## TOPIC\b")
BULLET_RE = re.compile(r"^(\s*)- (.+)$")
QA_ITEM_RE = re.compile(r"^\d+\.\s+\*\*Q:\*\*")


def _strip_inline_markdown(text: str) -> str:
    """Remove simple **bold** markers while keeping the label text."""
    return re.sub(r"\*\*(.+?)\*\*", r"\1", text).strip()


def _detect_format(path: Path) -> str:
    with path.open(encoding="utf-8") as handle:
        for raw_line in handle:
            if SLIDE_HEADING_RE.match(raw_line.strip()):
                return "m1"
    return "simple"


M1_FIELD_PREFIXES = (
    "**Title:**",
    "**Bullets:**",
    "**Speaker Notes:**",
    "**Q&A:**",
)


def parse_m1_txt_to_structure(filepath: str | Path) -> dict[str, Any]:
    """
    Parse Module-style structured text (e.g. M1.txt) into slides.

    Convention:
      # Document Title                    -> presentation metadata (not a slide)
      ## TOPIC X.X — Topic Name           -> topic header (not a slide)
      ### Slide N — Section Slide         -> section divider slide
      **Title:** Section name             -> title for the section slide above
      ### Slide N — Content Title         -> content slide
      **Bullets:**                        -> starts bullet list for current slide
      - bullet text                       -> slide bullet (supports **bold** labels)
      **Speaker Notes:** text             -> speaker notes for current slide
      **Q&A:**                            -> Q&A block appended to speaker notes
      1. **Q:** ... **A:** ...            -> individual Q&A items
    """
    structure: dict[str, Any] = {"title": "", "sections": []}
    current: dict[str, Any] | None = None
    mode: str | None = None
    path = Path(filepath)

    with path.open(encoding="utf-8") as handle:
        for raw_line in handle:
            line = raw_line.rstrip("\n\r")
            stripped = line.strip()

            if not stripped or stripped == "---":
                continue

            doc_title_match = DOC_TITLE_RE.match(stripped)
            if doc_title_match and not structure["title"]:
                structure["title"] = doc_title_match.group(1).strip()
                continue

            if TOPIC_HEADING_RE.match(stripped) or (
                stripped.startswith("**")
                and stripped.endswith("**")
                and not any(stripped.startswith(prefix) for prefix in M1_FIELD_PREFIXES)
            ):
                continue

            slide_match = SLIDE_HEADING_RE.match(stripped)
            if slide_match:
                slide_label = slide_match.group(1).strip()
                if slide_label.lower() == "section slide":
                    current = {
                        "type": "section_divider",
                        "title": "",
                        "bullets": [],
                        "notes": "",
                    }
                else:
                    current = {
                        "type": "content",
                        "title": slide_label,
                        "bullets": [],
                        "notes": "",
                    }
                structure["sections"].append(current)
                mode = None
                continue

            if current is None:
                logger.warning("Unrecognized line (no active slide): %r", stripped)
                continue

            if stripped.startswith("**Title:**"):
                current["title"] = _strip_inline_markdown(stripped[len("**Title:**") :])
                continue

            if stripped == "**Bullets:**":
                mode = "bullets"
                continue

            if stripped.startswith("**Speaker Notes:**"):
                note_text = _strip_inline_markdown(stripped[len("**Speaker Notes:**") :])
                current["notes"] = (
                    f"{current['notes']}\n\n{note_text}".strip()
                    if current["notes"]
                    else note_text
                )
                mode = None
                continue

            if stripped == "**Q&A:**":
                mode = "qa"
                if current["notes"]:
                    current["notes"] = f"{current['notes']}\n\nQ&A:"
                else:
                    current["notes"] = "Q&A:"
                continue

            if mode == "qa" and QA_ITEM_RE.match(stripped):
                qa_line = _strip_inline_markdown(stripped)
                current["notes"] = f"{current['notes']}\n{qa_line}"
                continue

            bullet_match = BULLET_RE.match(line)
            if bullet_match and mode in ("bullets", None):
                indent = len(bullet_match.group(1))
                level = 1 if indent >= 2 else 0
                current["bullets"].append(
                    {
                        "text": _strip_inline_markdown(bullet_match.group(2)),
                        "level": level,
                    }
                )
                mode = "bullets"
                continue

            logger.warning("Unrecognized line in slide %r: %r", current.get("title"), stripped)

    if not structure["sections"]:
        raise ValueError(f"No slides found in {path}. Expected '### Slide N — ...' headings.")

    for item in structure["sections"]:
        if item["type"] == "section_divider" and not item["title"]:
            raise ValueError(
                "Section slide is missing **Title:** immediately after '### Slide N — Section Slide'."
            )

    return structure


def parse_txt_to_structure(filepath: str | Path) -> dict[str, Any]:
    """Parse input text, auto-detecting simple or Module (M1) format."""
    path = Path(filepath)
    if _detect_format(path) == "m1":
        logger.info("Detected Module-style input format (### Slide headings)")
        return parse_m1_txt_to_structure(path)

    return _parse_simple_txt_to_structure(path)


def _parse_simple_txt_to_structure(filepath: Path) -> dict[str, Any]:
    """
    Parse a structured text file into an intermediate slide model.

    Convention:
      # Section Title       -> new section divider slide
      ## Slide Title        -> new content slide
      - bullet text         -> top-level bullet
        - sub bullet text   -> nested bullet (2-space indent before '-')
      Notes: ...            -> speaker notes for the current slide

    Blank lines are ignored. Lines before the first heading are ignored.
    """
    structure: dict[str, Any] = {"title": "", "sections": []}
    current: dict[str, Any] | None = None
    path = Path(filepath)

    with path.open(encoding="utf-8") as handle:
        for raw_line in handle:
            line = raw_line.rstrip("\n\r")
            if not line.strip():
                continue

            if line.startswith("# "):
                current = {
                    "type": "section_divider",
                    "title": line[2:].strip(),
                    "bullets": [],
                    "notes": "",
                }
                structure["sections"].append(current)
                if not structure["title"]:
                    structure["title"] = current["title"]
                continue

            if line.startswith("## "):
                current = {
                    "type": "content",
                    "title": line[3:].strip(),
                    "bullets": [],
                    "notes": "",
                }
                structure["sections"].append(current)
                continue

            if line.startswith("Notes:") and current is not None:
                note_text = line[len("Notes:") :].strip()
                if current["notes"]:
                    current["notes"] = f"{current['notes']}\n{note_text}"
                else:
                    current["notes"] = note_text
                continue

            bullet_match = re.match(r"^(\s*)- (.+)$", line)
            if bullet_match and current is not None:
                indent = len(bullet_match.group(1))
                level = 1 if indent >= 2 else 0
                current["bullets"].append(
                    {"text": bullet_match.group(2).strip(), "level": level}
                )
                continue

            logger.warning("Unrecognized line (no active slide): %r", line)

    if not structure["sections"]:
        raise ValueError(
            f"No slides found in {path}. Use '# Section' or '## Slide Title' headings."
        )

    return structure


def _placeholder_label(placeholder: Any) -> str:
    try:
        ph_type = placeholder.placeholder_format.type
        return PP_PLACEHOLDER(ph_type).name
    except (ValueError, AttributeError):
        return "UNKNOWN"


def validate_template(prs: Presentation) -> None:
    """Log layout names and placeholder indices so template mismatches fail fast."""
    logger.info("Template validation — %d slide layout(s):", len(prs.slide_layouts))
    for index, layout in enumerate(prs.slide_layouts):
        placeholders = list(layout.placeholders)
        if placeholders:
            details = ", ".join(
                f"idx={ph.placeholder_format.idx} ({_placeholder_label(ph)})"
                for ph in placeholders
            )
        else:
            details = "no placeholders"
        logger.info("  [%d] %r — %s", index, layout.name, details)


def find_layout(
    prs: Presentation,
    name_fragments: tuple[str, ...],
    fallback_index: int = 1,
) -> Any:
    """Return the first slide layout whose name contains any of the fragments."""
    lowered = tuple(fragment.lower() for fragment in name_fragments)
    for layout in prs.slide_layouts:
        layout_name = layout.name.lower()
        if any(fragment in layout_name for fragment in lowered):
            logger.debug("Matched layout %r for fragments %s", layout.name, name_fragments)
            return layout

    if fallback_index < len(prs.slide_layouts):
        fallback = prs.slide_layouts[fallback_index]
        logger.warning(
            "No layout matched %s; falling back to [%d] %r",
            name_fragments,
            fallback_index,
            fallback.name,
        )
        return fallback

    raise ValueError(
        f"Could not find a layout matching {name_fragments} "
        f"and fallback index {fallback_index} is out of range."
    )


def _find_title_placeholder(slide: Any) -> Any | None:
    if slide.shapes.title is not None:
        return slide.shapes.title
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            return placeholder
    return None


def _find_body_placeholder(slide: Any) -> Any | None:
    preferred_types = (
        PP_PLACEHOLDER.BODY,
        PP_PLACEHOLDER.OBJECT,
    )
    for ph_type in preferred_types:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == ph_type:
                return placeholder

    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx != 0:
            return placeholder
    return None


def _fill_bullets(text_frame: Any, bullets: list[dict[str, Any]]) -> None:
    text_frame.clear()
    first = True
    for bullet in bullets:
        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        paragraph.text = bullet["text"]
        paragraph.level = bullet["level"]
        first = False


def _clear_existing_slides(prs: Presentation) -> int:
    """Remove all slides from a presentation while keeping masters, layouts, and theme."""
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(slide_id)
    return len(slide_ids)


def build_pptx_from_structure(
    structure: dict[str, Any],
    template_path: str | Path,
    output_path: str | Path,
    *,
    clear_template_slides: bool = True,
) -> None:
    """Build a PPTX deck from parsed structure using template layouts and placeholders."""
    template_path = Path(template_path)
    output_path = Path(output_path)

    prs = Presentation(str(template_path))
    validate_template(prs)

    existing_slides = len(prs.slides)
    if existing_slides and clear_template_slides:
        removed = _clear_existing_slides(prs)
        logger.info(
            "Cleared %d existing slide(s) from template (theme/layouts preserved)",
            removed,
        )
    elif existing_slides:
        logger.warning(
            "Template contains %d existing slide(s); new slides will be appended",
            existing_slides,
        )

    section_layout = find_layout(prs, SECTION_LAYOUT_NAMES, fallback_index=0)
    content_layout = find_layout(prs, CONTENT_LAYOUT_NAMES, fallback_index=1)

    for slide_index, item in enumerate(structure["sections"], start=1):
        layout = (
            section_layout if item["type"] == "section_divider" else content_layout
        )
        slide = prs.slides.add_slide(layout)

        title_placeholder = _find_title_placeholder(slide)
        if title_placeholder is None:
            logger.error(
                "Slide %d (%r): no title placeholder on layout %r",
                slide_index,
                item["title"],
                layout.name,
            )
        else:
            title_placeholder.text = item["title"]
            logger.info(
                "Slide %d: type=%s, layout=%r, title placeholder idx=%s",
                slide_index,
                item["type"],
                layout.name,
                title_placeholder.placeholder_format.idx,
            )

        if item["bullets"]:
            body_placeholder = _find_body_placeholder(slide)
            if body_placeholder is None:
                logger.error(
                    "Slide %d (%r): no body placeholder for %d bullet(s) on layout %r",
                    slide_index,
                    item["title"],
                    len(item["bullets"]),
                    layout.name,
                )
            else:
                _fill_bullets(body_placeholder.text_frame, item["bullets"])
                logger.info(
                    "Slide %d: body placeholder idx=%s (%s), %d bullet(s)",
                    slide_index,
                    body_placeholder.placeholder_format.idx,
                    _placeholder_label(body_placeholder),
                    len(item["bullets"]),
                )

        if item["notes"]:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = item["notes"]
            logger.info("Slide %d: speaker notes set (%d chars)", slide_index, len(item["notes"]))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info("Saved %d slide(s) to %s", len(structure["sections"]), output_path)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Convert structured TXT/MD content into a theme-adaptive PPTX deck."
    )
    parser.add_argument("input_txt", help="Path to structured .txt or .md input file")
    parser.add_argument("template_pptx", help="Path to .pptx or .potx template file")
    parser.add_argument("output_pptx", help="Path for the generated .pptx output")
    parser.add_argument(
        "-v",
        "--verbose",
        action="store_true",
        help="Enable debug logging (layout matching, placeholder details)",
    )
    parser.add_argument(
        "--keep-template-slides",
        action="store_true",
        help="Keep existing slides in the template file (default: clear them first)",
    )
    args = parser.parse_args(argv)

    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(levelname)s: %(message)s",
    )

    structure = parse_txt_to_structure(args.input_txt)
    logger.info(
        "Parsed %d slide(s) from %s (presentation title: %r)",
        len(structure["sections"]),
        args.input_txt,
        structure["title"] or "(none)",
    )

    build_pptx_from_structure(
        structure,
        args.template_pptx,
        args.output_pptx,
        clear_template_slides=not args.keep_template_slides,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
